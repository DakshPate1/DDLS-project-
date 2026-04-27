"""Build the 12-slide DDLS interim presentation as a .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy, os

# ── palette ──────────────────────────────────────────────────────────────────
C_BG      = RGBColor(0x0D, 0x1B, 0x2A)   # dark navy
C_ACCENT  = RGBColor(0x1E, 0x90, 0xFF)   # dodger blue
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xCC, 0xDD, 0xEE)
C_GREEN   = RGBColor(0x3C, 0xD6, 0x8A)
C_YELLOW  = RGBColor(0xFF, 0xD7, 0x00)
C_RED     = RGBColor(0xFF, 0x4C, 0x4C)
C_GREY    = RGBColor(0x88, 0x99, 0xAA)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank

# ── helpers ──────────────────────────────────────────────────────────────────

def add_slide():
    sl = prs.slides.add_slide(BLANK)
    bg = sl.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C_BG
    return sl


def tb(slide, text, l, t, w, h, size=20, bold=False, color=C_WHITE,
       align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def rule(slide, l, t, w, h=2, color=C_ACCENT):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        l, t, w, Pt(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def slide_header(slide, title, subtitle=None):
    rule(slide, Inches(0.4), Inches(0.35), Inches(12.5))
    tb(slide, title, Inches(0.4), Inches(0.45), Inches(12.5), Inches(0.8),
       size=28, bold=True, color=C_WHITE)
    if subtitle:
        tb(slide, subtitle, Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.5),
           size=16, color=C_LIGHT)


def bullet_block(slide, items, l, t, w, h, size=18, marker="•", color=C_WHITE,
                 line_spacing=1.2):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_before = Pt(4)
        run = para.add_run()
        bullet_char = marker + "  " if marker else ""
        run.text = bullet_char + item
        run.font.size = Pt(size)
        run.font.color.rgb = color


def colored_box(slide, l, t, w, h, fill_color, text=None, text_size=18,
                text_color=C_WHITE, bold=False):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.size = Pt(text_size)
        run.font.color.rgb = text_color
        run.font.bold = bold
    return shape


def table_on_slide(slide, headers, rows, l, t, w, col_widths=None,
                   header_color=C_ACCENT, row_colors=None,
                   font_size=15, header_font_size=15):
    n_cols = len(headers)
    n_rows = len(rows)
    row_h = Inches(0.42)
    total_h = row_h * (n_rows + 1)
    tbl = slide.shapes.add_table(n_rows + 1, n_cols, l, t, w, total_h).table
    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = cw
    # header
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.bold = True
        run.font.size = Pt(header_font_size)
        run.font.color.rgb = C_WHITE
    # rows
    for ri, row in enumerate(rows):
        bg = row_colors[ri] if row_colors and ri < len(row_colors) else RGBColor(0x12, 0x26, 0x3A)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(font_size)
            run.font.color.rgb = C_WHITE
    return tbl


# ── image helper ─────────────────────────────────────────────────────────────

PLOTS = "/home/ahmad/unibe/DDLS-project-/plots"

def add_image(slide, fname, l, t, w, h):
    path = os.path.join(PLOTS, fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        colored_box(slide, l, t, w, h, RGBColor(0x22, 0x33, 0x44),
                    text=f"[{fname}]", text_size=14, text_color=C_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rule(sl, Inches(0.5), Inches(2.85), Inches(12.3), h=3, color=C_ACCENT)
tb(sl, "Federated Risk-Sensitive Q-Learning\nin Continuous Time",
   Inches(0.5), Inches(1.05), Inches(12.3), Inches(1.8),
   size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
tb(sl, "Reproducing Xie (2025) and extending to heterogeneous financial markets",
   Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.6),
   size=20, color=C_LIGHT, align=PP_ALIGN.CENTER, italic=True)
tb(sl, "DDLS Spring 2026  ·  Universität Bern  ·  April 28, 2026",
   Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.5),
   size=16, color=C_GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Why This Paper?
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
slide_header(sl, "Why This Paper?", "Finance is where risk-sensitive RL actually matters")

# left box — The Problem
colored_box(sl, Inches(0.4), Inches(1.7), Inches(5.9), Inches(0.45),
            RGBColor(0x1A, 0x30, 0x4A), text="THE PROBLEM", text_size=14,
            text_color=C_ACCENT, bold=True)
bullet_block(sl, [
    "Portfolio allocation is a continuous-time decision",
    "Banks, pension funds, hedge funds all face it daily",
    "Classical RL: risk-neutral agents + discrete time",
    "Neither assumption holds in real finance",
], Inches(0.4), Inches(2.2), Inches(5.9), Inches(2.8), size=17)

# right box — The Fix
colored_box(sl, Inches(7.0), Inches(1.7), Inches(5.9), Inches(0.45),
            RGBColor(0x0A, 0x3A, 0x2A), text="THE FIX", text_size=14,
            text_color=C_GREEN, bold=True)
bullet_block(sl, [
    "CT-RS-q: continuous-time risk-sensitive q-learning",
    "Handles non-linear objectives: mean-variance, CVaR,\n   exponential utility",
    "Risk-sensitivity baked into the algorithm via OCE\n   framework",
    "Provably optimal policy",
], Inches(7.0), Inches(2.2), Inches(5.9), Inches(2.8), size=17, color=C_WHITE)

tb(sl, '"Most RL papers either ignore risk or bolt it on as a constraint.\nThis paper bakes risk-sensitivity into the learning algorithm itself."',
   Inches(0.4), Inches(5.2), Inches(12.5), Inches(0.9),
   size=15, color=C_YELLOW, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Problem Setup
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
slide_header(sl, "The Problem Setup", "Wealth follows a stochastic process; maximize a risk-sensitive functional")

# SDE box
colored_box(sl, Inches(0.4), Inches(1.65), Inches(12.5), Inches(0.38),
            RGBColor(0x10, 0x25, 0x40), text="Portfolio SDE", text_size=13,
            text_color=C_ACCENT, bold=True)
tb(sl, "dX_t = X_t (a_t·r₁ + (1−a_t)·r₂) dt  +  X_t (a_t·σ₁·dW₁ + (1−a_t)·σ₂·dW₂)",
   Inches(0.6), Inches(2.1), Inches(12.0), Inches(0.55),
   size=17, color=C_WHITE, bold=True)
tb(sl, "a_t ∈ [0,1]  — fraction of wealth in asset 1 at time t",
   Inches(0.6), Inches(2.65), Inches(10.0), Inches(0.4),
   size=15, color=C_LIGHT)

# Objective
colored_box(sl, Inches(0.4), Inches(3.15), Inches(12.5), Inches(0.38),
            RGBColor(0x10, 0x25, 0x40), text="Objective (Mean-Variance)", text_size=13,
            text_color=C_ACCENT, bold=True)
tb(sl, "MV(X_T)  =  E[X_T]  −  (α/2)·Var(X_T)         α = risk aversion",
   Inches(0.6), Inches(3.6), Inches(12.0), Inches(0.5),
   size=17, color=C_WHITE, bold=True)

# Why hard + fix
colored_box(sl, Inches(0.4), Inches(4.2), Inches(5.9), Inches(0.38),
            RGBColor(0x3A, 0x10, 0x10), text="WHY CLASSICAL BELLMAN FAILS", text_size=13,
            text_color=C_RED, bold=True)
tb(sl, "Non-linear objectives break the Markov property.\nOptimal policy is not Markovian in just (t, X_t).",
   Inches(0.4), Inches(4.65), Inches(5.9), Inches(1.0), size=16, color=C_WHITE)

colored_box(sl, Inches(6.6), Inches(4.2), Inches(6.3), Inches(0.38),
            RGBColor(0x0A, 0x35, 0x20), text="THE FIX — STATE AUGMENTATION", text_size=13,
            text_color=C_GREEN, bold=True)
tb(sl, "Augment state to  (t, X_t, B⁰_t, B¹_t)\n"
       "B⁰ accumulates rewards; B¹ tracks OCE discounting\n"
       "→ Bellman works again on the augmented state",
   Inches(6.6), Inches(4.65), Inches(6.3), Inches(1.1), size=16, color=C_WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — The Algorithm
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
slide_header(sl, "The Algorithm: CT-RS-q", "Jointly learns value function J_θ and q-function q_ψ via continuous-time TD")

# Two objects
colored_box(sl, Inches(0.4), Inches(1.65), Inches(6.0), Inches(0.4),
            C_ACCENT, text="Two Learned Objects  (8 scalar parameters total)", text_size=14,
            text_color=C_WHITE, bold=True)
bullet_block(sl, [
    "J_θ(t,x) — value function, quadratic in x  [3 params: θ]",
    "q_ψ(t,x,a) — q-function, quadratic in (a−a*)  [5 params: ψ]",
    "No neural nets — provably optimal parameterization",
], Inches(0.4), Inches(2.12), Inches(6.0), Inches(1.3), size=16)

# Policy
colored_box(sl, Inches(6.7), Inches(1.65), Inches(6.2), Inches(0.4),
            RGBColor(0x0A, 0x35, 0x20), text="Policy (Gaussian)", text_size=14,
            text_color=C_GREEN, bold=True)
tb(sl, "π_ψ(a | t, x)  =  N( a*(t,x),  σ²_τ(t,x) )\n\n"
       "mean = optimal allocation   a*(t,x)\n"
       "variance = exploration temperature  τ",
   Inches(6.7), Inches(2.12), Inches(6.2), Inches(1.3), size=16, color=C_WHITE)

# Loop boxes
loop_items = [
    ("1. ROLLOUT", "Run trajectory under current policy π_ψ", C_ACCENT),
    ("2. TD ERRORS", "TD_k = J_{k+1} − J_k − q_k·dt", RGBColor(0x1A, 0x60, 0x30)),
    ("3. GRADIENT UPDATE", "Ascent on θ and ψ using TD as multiplier", RGBColor(0x50, 0x20, 0x80)),
    ("4. UPDATED POLICY", "New π_ψ with refined mean & variance", RGBColor(0x80, 0x40, 0x00)),
]
box_w = Inches(2.9)
for i, (title, desc, col) in enumerate(loop_items):
    lx = Inches(0.4) + i * (box_w + Inches(0.2))
    colored_box(sl, lx, Inches(3.7), box_w, Inches(0.35),
                col, text=title, text_size=13, text_color=C_WHITE, bold=True)
    tb(sl, desc, lx, Inches(4.1), box_w, Inches(0.65), size=14, color=C_LIGHT)
    if i < 3:
        tb(sl, "→", lx + box_w + Inches(0.01), Inches(3.85), Inches(0.2), Inches(0.4),
           size=22, color=C_ACCENT, bold=True)

tb(sl, "\"The learned object is an allocation rule — at every instant t, given current wealth x,\n"
       "it tells you what fraction to invest in asset 1.\"",
   Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.8),
   size=15, color=C_YELLOW, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — What We Reproduced
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
slide_header(sl, "What We Reproduced", "Table 1, Figure 1, Figure 2 — MV objective within 1.5% of the paper")

headers = ["Policy", "Cum Return", "Std Dev", "MV Objective"]
rows = [
    ["Baseline (a=0.5)  ours / paper",  "0.221 / 0.222",  "0.095 / 0.096",  "1.217 / 1.217"],
    ["Optimal (closed form)  ours / paper", "0.709 / 0.713", "0.720 / 0.721", "1.450 / 1.453"],
    ["CT-RS-q (learned)  ours / paper",  "0.729 / 0.816",  "0.791 / 0.872",  "1.416 / 1.437"],
]
row_colors = [
    RGBColor(0x10, 0x25, 0x40),
    RGBColor(0x10, 0x25, 0x40),
    RGBColor(0x0A, 0x35, 0x20),
]
col_widths = [Inches(5.2), Inches(2.3), Inches(2.3), Inches(2.3)]
table_on_slide(sl, headers, rows, Inches(0.4), Inches(1.7), Inches(12.5),
               col_widths=col_widths, row_colors=row_colors)

bullet_block(sl, [
    "Baseline and Optimal land within Monte Carlo noise — essentially exact",
    "CT-RS-q: MV within 1.5% of paper — the 10.7% cumulative return gap traces to a single",
    "   parameter overshooting on this seed (stochasticity, not a code error)",
], Inches(0.4), Inches(4.55), Inches(12.5), Inches(1.4), size=17, color=C_LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Figure 1 and Figure 2
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
slide_header(sl, "Convergence & Time Evolution", "Qualitative picture matches the paper — two findings to highlight")

add_image(sl, "figure1_convergence.png",
          Inches(0.3), Inches(1.5), Inches(6.3), Inches(5.6))
add_image(sl, "figure2_time_evolution.png",
          Inches(6.8), Inches(1.5), Inches(6.1), Inches(5.6))

tb(sl, "Fig 1 — 8-panel parameter convergence\nNote ψ_ce1 and ψ_ce2 drifting in parallel →",
   Inches(0.3), Inches(6.6), Inches(6.3), Inches(0.7),
   size=13, color=C_YELLOW, italic=True)
tb(sl, "Fig 2 — CT-RS-q tracks optimal on MV\nBoth strongly beat baseline",
   Inches(6.8), Inches(6.6), Inches(6.1), Inches(0.7),
   size=13, color=C_GREEN, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — The Key Finding (structural non-identifiability)
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
slide_header(sl, "What We Found: Structural Non-Identifiability",
             "The paper's parameterization has a flat direction the authors don't acknowledge")

# Observation
colored_box(sl, Inches(0.4), Inches(1.6), Inches(12.5), Inches(0.38),
            RGBColor(0x1A, 0x10, 0x30), text="The Observation", text_size=14,
            text_color=C_ACCENT, bold=True)
tb(sl, "ψ_ce1 and ψ_ce2 drift in parallel in Figure 1.  "
       "Individual values diverge from ground truth — but their DIFFERENCE stays constant.",
   Inches(0.4), Inches(2.05), Inches(12.5), Inches(0.55), size=16, color=C_WHITE)

# Math
colored_box(sl, Inches(0.4), Inches(2.7), Inches(12.5), Inches(0.38),
            RGBColor(0x10, 0x10, 0x30), text="The Math", text_size=14,
            text_color=C_ACCENT, bold=True)
tb(sl, "a_ψ  ∝  ψ_a0 − ψ_a1·(1 + c1_ψ / (2·c2_ψ·x))        "
       "c1_ψ / c2_ψ  ∝  exp( (ψ_ce1 − ψ_ce2)·τ )",
   Inches(0.4), Inches(3.15), Inches(12.5), Inches(0.5), size=16, color=C_YELLOW, bold=True)
tb(sl, "ψ_ce1 and ψ_ce2 enter the loss ONLY through their difference.\n"
       "Gradient is zero along (+1, +1) direction  →  flat null-space.  No training resolves individual values.",
   Inches(0.4), Inches(3.7), Inches(12.5), Inches(0.65), size=15, color=C_LIGHT)

# What the paper says vs reality
colored_box(sl, Inches(0.4), Inches(4.45), Inches(5.8), Inches(0.38),
            RGBColor(0x3A, 0x10, 0x10), text="Paper says", text_size=13,
            text_color=C_RED, bold=True)
tb(sl, '"ψ_ce1 and ψ_ce2 do not fully converge due to non-zero τ"\n'
       '→  Incomplete: cause is structural, not an exploration artifact',
   Inches(0.4), Inches(4.9), Inches(5.8), Inches(0.95), size=15, color=C_WHITE)

# Implication for FL
colored_box(sl, Inches(6.7), Inches(4.45), Inches(6.2), Inches(0.38),
            RGBColor(0x0A, 0x35, 0x20), text="Why FedAvg Fails Here", text_size=13,
            text_color=C_GREEN, bold=True)
tb(sl, "Each worker drifts to a DIFFERENT point in the null-space.\n"
       "FedAvg averages those points — a third meaningless point.\n"
       "→  FedAvg is broken by construction on this loss.",
   Inches(6.7), Inches(4.9), Inches(6.2), Inches(0.95), size=15, color=C_WHITE)

tb(sl, "We found something the authors missed — and it turns out to be the exact reason vanilla FedAvg fails.",
   Inches(0.4), Inches(6.05), Inches(12.5), Inches(0.55),
   size=15, color=C_YELLOW, italic=True, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — The Federated Setting
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
slide_header(sl, "The Federated Setting",
             "Finance is naturally federated — this is not a forced extension")

# motivation
bullet_block(sl, [
    "Banks, asset managers, pension funds: each has its own market data, its own clients, its own risk tolerance",
    "They cannot and do not share raw data",
    "Each institution has a different risk aversion α — a pension fund and a hedge fund should NOT share a policy",
], Inches(0.4), Inches(1.6), Inches(12.5), Inches(1.5), size=18)

# worker table
headers2 = ["Worker", "r₁, r₂", "σ₁, σ₂", "Risk α"]
rows2 = [
    ["Worker 1", "0.15,  0.25", "0.10,  0.12", "1.0"],
    ["Worker 2", "0.30,  0.10", "0.20,  0.15", "2.0"],
    ["Worker 3", "0.05,  0.08", "0.05,  0.06", "0.5"],
    ["Worker 4", "0.20,  0.20", "0.25,  0.30", "3.0"],
    ["Held-out market", "0.12,  0.18", "0.15,  0.18", "1.5"],
]
row_colors2 = [
    RGBColor(0x10, 0x25, 0x40),
    RGBColor(0x12, 0x28, 0x44),
    RGBColor(0x10, 0x25, 0x40),
    RGBColor(0x12, 0x28, 0x44),
    RGBColor(0x0A, 0x30, 0x1A),
]
col_widths2 = [Inches(2.8), Inches(3.0), Inches(3.0), Inches(2.0)]
table_on_slide(sl, headers2, rows2, Inches(1.5), Inches(3.2), Inches(10.5),
               col_widths=col_widths2, row_colors=row_colors2, font_size=16)

tb(sl, "Heterogeneous markets AND heterogeneous risk preferences — both dimensions the original paper has no mechanism for.",
   Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.6),
   size=15, color=C_YELLOW, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Our Extension: Three Methods
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
slide_header(sl, "Our Extension: Three Methods",
             "Local (no comms) vs FedAvg (fails) vs PF-CT-RS-q (our contribution)")

method_data = [
    ("LOCAL CT-RS-q", "Each worker trains independently.\nNo communication.\nBaseline — upper bound for personalization.",
     C_GREY, "✓ Personalized\n✗ No knowledge sharing"),
    ("VANILLA FedAvg", "Server averages ALL parameters (θ and ψ).\nFails because:\n"
     "1. Null-space averaging makes ψ_ce1, ψ_ce2 meaningless\n"
     "2. Averaging ψ across different α destroys risk calibration",
     C_RED, "✗ Null-space collapse\n✗ Risk calibration lost"),
    ("PF-CT-RS-q  (ours)", "Federate θ only (shared market structure)\nKeep ψ local (risk preferences are personal)\n"
     "Null-space fix: freeze ψ_ce2 — pins the flat direction\n"
     "θ / ψ split comes directly from the math of the paper",
     C_GREEN, "✓ Shared dynamics\n✓ Personalized risk\n✓ Null-space pinned"),
]

col_w = Inches(4.1)
for i, (title, desc, col, verdict) in enumerate(method_data):
    lx = Inches(0.3) + i * (col_w + Inches(0.2))
    colored_box(sl, lx, Inches(1.6), col_w, Inches(0.4),
                col, text=title, text_size=14, text_color=C_WHITE, bold=True)
    tb(sl, desc, lx, Inches(2.1), col_w, Inches(2.6), size=14, color=C_LIGHT)
    colored_box(sl, lx, Inches(4.8), col_w, Inches(1.0),
                RGBColor(0x10, 0x22, 0x38),
                text=verdict, text_size=14, text_color=col, bold=False)

tb(sl, "The θ/ψ split is principled: θ captures dynamics (transferable), ψ captures risk preferences (personal).",
   Inches(0.3), Inches(6.1), Inches(12.5), Inches(0.6),
   size=15, color=C_YELLOW, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Decentralized Extension
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
slide_header(sl, "Decentralized Extension",
             "Same principled split, no server needed — connects directly to Week 7 (Dec-RL)")

# Setup
colored_box(sl, Inches(0.4), Inches(1.65), Inches(5.9), Inches(0.4),
            RGBColor(0x10, 0x25, 0x40), text="Setup", text_size=14,
            text_color=C_ACCENT, bold=True)
bullet_block(sl, [
    "Remove central server — workers communicate peer-to-peer",
    "Two topologies: Ring (2 neighbors)  and  Fully Connected",
    "Mixing weights: Metropolis-Hastings (standard for gossip)",
    "Same parameter split: gossip θ, keep ψ local",
    "Same null-space fix: freeze ψ_ce2",
], Inches(0.4), Inches(2.1), Inches(5.9), Inches(2.6), size=16)

# Why interesting
colored_box(sl, Inches(6.7), Inches(1.65), Inches(6.2), Inches(0.4),
            RGBColor(0x0A, 0x35, 0x20), text="Why Harder Than Federated", text_size=14,
            text_color=C_GREEN, bold=True)
bullet_block(sl, [
    "No central aggregator to correct null-space drift",
    "Each worker's θ must converge via local averaging only",
    "Freeze on ψ_ce2 becomes even more critical",
    "FC topology: faster consensus  (O(log n) mixing)",
    "Ring topology: slower consensus — tests robustness",
], Inches(6.7), Inches(2.1), Inches(6.2), Inches(2.6), size=16, color=C_WHITE)

# Topology diagrams (text-based since we can't draw circles)
colored_box(sl, Inches(0.4), Inches(4.95), Inches(5.9), Inches(1.9),
            RGBColor(0x0D, 0x1E, 0x35),
            text="Ring Topology\n1 — 2 — 3 — 4 — 1\n(2 neighbors each)",
            text_size=16, text_color=C_LIGHT)
colored_box(sl, Inches(6.7), Inches(4.95), Inches(6.2), Inches(1.9),
            RGBColor(0x0D, 0x1E, 0x35),
            text="Fully Connected\nEvery worker ↔ every other\n(max connectivity, fastest mixing)",
            text_size=16, text_color=C_LIGHT)

tb(sl, "Phase 3 — in design. Building after Phase 2 is verified (target: May 19).",
   Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.3),
   size=14, color=C_GREY, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Results Plan
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
slide_header(sl, "What We Will Show", "Concrete, honest plan — one primary plot, honest timeline")

# Expected results block
colored_box(sl, Inches(0.4), Inches(1.6), Inches(7.5), Inches(0.4),
            C_ACCENT, text="Primary Result: MV on held-out market", text_size=14,
            text_color=C_WHITE, bold=True)

result_rows = [
    ("Local CT-RS-q  (avg over workers)", "─────────────────────", "X (reference)"),
    ("Vanilla FedAvg", "────────── X", "drops — null-space collapse"),
    ("PF-CT-RS-q  (ours)", "─────────────────────", "✓  principal result"),
    ("Dec CT-RS-q  ring", "───────────────────", "✓  slightly lower"),
    ("Dec CT-RS-q  FC", "─────────────────────", "✓  matches federated"),
]
for i, (method, bar, note) in enumerate(result_rows):
    t_y = Inches(2.1) + i * Inches(0.58)
    col = C_GREEN if "✓" in note else (C_RED if "drops" in note else C_GREY)
    tb(sl, method, Inches(0.5), t_y, Inches(4.2), Inches(0.5), size=15, color=C_LIGHT)
    tb(sl, bar, Inches(4.8), t_y, Inches(3.0), Inches(0.5), size=15, color=col)
    tb(sl, note, Inches(8.0), t_y, Inches(4.8), Inches(0.5), size=14, color=col)

# Timeline
colored_box(sl, Inches(8.2), Inches(1.6), Inches(4.8), Inches(0.4),
            RGBColor(0x10, 0x25, 0x40), text="Timeline", text_size=14,
            text_color=C_ACCENT, bold=True)
timeline = [
    ("Phase 2 (federated)", "May 12"),
    ("Phase 3 (decentralized)", "May 19"),
    ("Final report draft", "May 22"),
    ("Final presentation", "May 26"),
]
for i, (item, date) in enumerate(timeline):
    t_y = Inches(2.1) + i * Inches(0.65)
    tb(sl, item, Inches(8.2), t_y, Inches(3.4), Inches(0.55), size=15, color=C_LIGHT)
    tb(sl, date, Inches(11.8), t_y, Inches(1.2), Inches(0.55), size=15, color=C_YELLOW, bold=True)

tb(sl, '"We\'re not just applying FedAvg — we\'re using a property of the paper\'s own parameterization\n'
       'to motivate why it will fail and how to fix it. That\'s the scientific contribution."',
   Inches(0.4), Inches(5.9), Inches(12.5), Inches(0.9),
   size=15, color=C_YELLOW, italic=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Summary
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
slide_header(sl, "Summary", "Phase 1 done, extension principled and locked, ready for Phase 2")

items = [
    ("✅", "Reproduced Table 1, Figure 1, Figure 2  (MV within 1.5% of paper)", C_GREEN),
    ("✅", "Found and formalized structural non-identifiability in B.2 parameterization", C_GREEN),
    ("✅", "Extension design: θ federation, ψ personalization, null-space pinning", C_GREEN),
    ("⏳", "Phase 2 (Federated, Module FL): in progress — target May 12", C_YELLOW),
    ("⏳", "Phase 3 (Decentralized RL, Module Dec-RL): after Phase 2 verified — May 19", C_YELLOW),
]
for i, (icon, text, col) in enumerate(items):
    t_y = Inches(1.75) + i * Inches(0.82)
    colored_box(sl, Inches(0.4), t_y, Inches(0.55), Inches(0.55),
                RGBColor(0x10, 0x25, 0x40), text=icon, text_size=20,
                text_color=col, bold=True)
    tb(sl, text, Inches(1.1), t_y + Inches(0.05), Inches(11.5), Inches(0.6),
       size=19, color=col if col == C_YELLOW else C_WHITE)

rule(sl, Inches(0.4), Inches(6.2), Inches(12.5), h=2, color=C_ACCENT)
tb(sl, "This paper is a strong choice for DDLS extensions because the math of continuous-time\n"
       "risk-sensitive RL gives you a principled reason why federated averaging fails — and a principled fix.",
   Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.9),
   size=17, color=C_LIGHT, align=PP_ALIGN.CENTER, italic=True)


# ── save ─────────────────────────────────────────────────────────────────────
out_path = "/home/ahmad/unibe/DDLS-project-/DDLS_Presentation_April28.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
